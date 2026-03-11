import streamlit as st
import ollama
import uuid
import json
import hashlib
import os
from datetime import datetime
from pathlib import Path
from docx import Document
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor

# --- CONFIGURATION ---
st.set_page_config(
    layout="wide", page_title="Narrative Sanctuary", page_icon="🌌")
CHAT_FILE  = Path("chat_history.json")
USERS_FILE = Path("data/stories.json")
GENRES = ["Fantasy", "Sci-Fi", "Thriller", "Romance", "Horror", "Adventure", "Humorous", "Historical"]
MODES = ["Fast", "Thinking", "Creative"]

# --- IMAGE BACKGROUNDS ---
BACKGROUND_IMAGES = {
    "Default": "https://images.unsplash.com/photo-1451187580459-43490279c0fa?q=80&w=2000",
    "Fantasy": "https://i.pinimg.com/736x/62/08/10/6208106d31891e449b1e3feb622e0f0e.jpg",
    "Sci-Fi": "https://i.pinimg.com/1200x/cb/0f/f0/cb0ff0cd14f228b38cb31e6ef35200d2.jpg",
    "Thriller": "https://i.pinimg.com/1200x/bd/41/15/bd41150d60df6e4fe55564fba654f266.jpg",
    "Romance": "https://i.pinimg.com/736x/b5/1a/19/b51a19f60575e934aef6b0178f3982d5.jpg",
    "Horror": "https://images.unsplash.com/photo-1509248961158-e54f6934749c?q=80&w=2000",
    "Adventure": "https://images.unsplash.com/photo-1464822759023-fed622ff2c3b?q=80&w=2000",
    "Humorous": "https://images.unsplash.com/photo-1527224538127-2104bb71c51b?q=80&w=2000",
    "Historical": "https://images.unsplash.com/photo-1524492412937-b28074a5d7da?q=80&w=2000",
}

# --- DATA ENGINE ---


def load_data():
    if CHAT_FILE.exists():
        try:
            with open(CHAT_FILE, "r") as f:
                return json.load(f)
        except:
            return {}
    return {}


def save_data():
    with open(CHAT_FILE, "w") as f:
        json.dump(st.session_state.conversations, f, indent=4)


# --- USER AUTH HELPERS ---
def load_users():
    USERS_FILE.parent.mkdir(parents=True, exist_ok=True)
    if USERS_FILE.exists():
        try:
            with open(USERS_FILE, "r") as f:
                return json.load(f)
        except:
            pass
    return {"users": {}}

def save_users(data):
    USERS_FILE.parent.mkdir(parents=True, exist_ok=True)
    with open(USERS_FILE, "w") as f:
        json.dump(data, f, indent=2)

def hash_password(password, salt=None):
    if salt is None:
        salt = os.urandom(16).hex()
    pw_hash = hashlib.sha256((salt + password).encode()).hexdigest()
    return pw_hash, salt


# --- LOGIN / REGISTER PAGE ---
def show_auth_page():
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Space+Grotesk:wght@300;400;600;700&display=swap');
    html, body, [data-testid="stApp"] { font-family: 'Space Grotesk', sans-serif; }

    [data-testid="stApp"] {
        background: linear-gradient(180deg,#0a1520 0%,#112235 50%,#0d1b2a 100%) !important;
        min-height: 100vh;
    }
    [data-testid="stApp"]::before {
        content:"";
        position:fixed; inset:0; z-index:0;
        background:
          radial-gradient(ellipse 80% 50% at 50% 0%, #1a3a5c33 0%, transparent 70%),
          url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 1440 320' preserveAspectRatio='none'%3E%3Cpath fill='%230c1e30' d='M0,300L20,180L40,220L60,110L80,160L110,130L140,80L170,130L200,110L230,60L260,110L300,90L340,40L380,90L420,70L460,20L500,70L540,50L580,10L620,50L660,40L700,0L740,40L780,30L820,0L860,30L900,20L940,0L980,20L1020,10L1060,0L1100,10L1140,20L1180,0L1220,20L1260,30L1300,10L1340,30L1380,60L1440,45L1440,320Z'/%3E%3C/svg%3E")
          bottom/100% auto no-repeat;
    }
    [data-testid="stHeader"], footer { display:none !important; }

    /* card */
    div[data-testid="stForm"] {
        background: rgba(180,205,230,0.13) !important;
        backdrop-filter: blur(18px) saturate(1.4) !important;
        border: 1px solid rgba(255,255,255,0.22) !important;
        border-radius: 20px !important;
        padding: 10px 24px 24px !important;
        box-shadow: 0 8px 40px rgba(0,0,0,0.45) !important;
    }
    /* inputs */
    input[type="text"], input[type="password"], input[type="email"] {
        background: rgba(255,255,255,0.1) !important;
        border: 1px solid rgba(255,255,255,0.2) !important;
        border-radius: 10px !important;
        color: #e0eaf5 !important;
    }
    input[type="text"]:focus, input[type="password"]:focus {
        border-color: rgba(120,180,240,0.7) !important;
        background: rgba(255,255,255,0.17) !important;
    }
    /* submit button */
    [data-testid="stFormSubmitButton"] > button {
        width:100% !important;
        background: #1a2e42 !important;
        border-radius: 10px !important;
        color: #e0eaf5 !important;
        font-weight: 600 !important;
        border: none !important;
        transition: background 0.2s !important;
    }
    [data-testid="stFormSubmitButton"] > button:hover { background:#243d56 !important; }
    </style>
    """, unsafe_allow_html=True)

    # Tabs: Login / Register
    tab_login, tab_reg = st.tabs(["🔑  Login", "📝  Register"])
    db = load_users()

    # ── Login Tab ──
    with tab_login:
        st.markdown("<h3 style='text-align:center;color:#e0eaf5;margin-bottom:8px;'>Welcome Back</h3>", unsafe_allow_html=True)
        with st.form("login_form"):
            email    = st.text_input("Email",    placeholder="your@email.com")
            password = st.text_input("Password", placeholder="Password", type="password")
            submitted = st.form_submit_button("Login")

        if submitted:
            users = db.get("users", {})
            user  = users.get(email)
            if user:
                pw_hash, _ = hash_password(password, user["salt"])
                if pw_hash == user["password_hash"]:
                    st.session_state.logged_in    = True
                    st.session_state.current_user = email
                    st.success("Login successful! Loading your stories…")
                    st.rerun()
                else:
                    st.error("Incorrect password.")
            else:
                st.error("No account found with that email.")

        st.markdown("<p style='text-align:center;color:#8ab;font-size:0.82rem;margin-top:8px;'>Forgot password? Contact support.</p>", unsafe_allow_html=True)

    # ── Register Tab ──
    with tab_reg:
        st.markdown("<h3 style='text-align:center;color:#e0eaf5;margin-bottom:8px;'>Create Account</h3>", unsafe_allow_html=True)
        with st.form("register_form"):
            r_name     = st.text_input("Full Name",        placeholder="Your name")
            r_email    = st.text_input("Email",            placeholder="your@email.com")
            r_password = st.text_input("Password",         placeholder="Min 6 characters", type="password")
            r_confirm  = st.text_input("Confirm Password", placeholder="Repeat password",   type="password")
            reg_submit = st.form_submit_button("Create Account")

        if reg_submit:
            if len(r_password) < 6:
                st.error("Password must be at least 6 characters.")
            elif r_password != r_confirm:
                st.error("Passwords do not match.")
            elif not r_email or not r_name:
                st.error("Please fill in all fields.")
            elif r_email in db.get("users", {}):
                st.error("An account with this email already exists.")
            else:
                pw_hash, salt = hash_password(r_password)
                db.setdefault("users", {})[r_email] = {
                    "name":          r_name,
                    "password_hash": pw_hash,
                    "salt":          salt,
                    "created_at":    datetime.now().isoformat(),
                    "stories":       []
                }
                save_users(db)
                st.success(f"Account created for {r_name}! You can now login.")


def generate_docx(data):
    doc = Document()
    doc.add_heading(f"Manuscript: {data['genre']}", 0)
    doc.add_paragraph(f"Author: {data['profile']['name']}")
    doc.add_paragraph(f"Synopsis: {data['profile']['bio']}")
    doc.add_paragraph("_" * 20)
    for msg in data['messages']:
        role = data['profile']['name'] if msg['role'] == "user" else "Narrator"
        p = doc.add_paragraph()
        p.add_run(f"{role}: ").bold = (msg['role'] == "user")
        p.add_run(msg['content']).italic = (msg['role'] != "user")
    bio = BytesIO()
    doc.save(bio)
    return bio.getvalue()

def generate_pptx(data):
    GENRE_THEMES = {
        "Fantasy":  {"bg": RGBColor(0x1a,0x00,0x33), "accent": RGBColor(0x9b,0x59,0xb6), "title": RGBColor(0xff,0xd7,0x00)},
        "Sci-Fi":   {"bg": RGBColor(0x00,0x11,0x33), "accent": RGBColor(0x00,0xbc,0xd4), "title": RGBColor(0x00,0xff,0xff)},
        "Thriller": {"bg": RGBColor(0x1a,0x1a,0x1a), "accent": RGBColor(0xe7,0x4c,0x3c), "title": RGBColor(0xff,0x44,0x44)},
        "Romance":  {"bg": RGBColor(0x2d,0x00,0x18), "accent": RGBColor(0xe9,0x1e,0x8c), "title": RGBColor(0xff,0xb6,0xc1)},
        "Horror":      {"bg": RGBColor(0x0d,0x00,0x00), "accent": RGBColor(0x8b,0x00,0x00), "title": RGBColor(0xcc,0x00,0x00)},
        "Adventure":   {"bg": RGBColor(0x0d,0x1f,0x0a), "accent": RGBColor(0x2e,0x8b,0x57), "title": RGBColor(0x7c,0xfc,0x00)},
        "Humorous":    {"bg": RGBColor(0x1a,0x1a,0x00), "accent": RGBColor(0xff,0xd7,0x00), "title": RGBColor(0xff,0xa5,0x00)},
        "Historical":  {"bg": RGBColor(0x1a,0x10,0x00), "accent": RGBColor(0x8b,0x6f,0x47), "title": RGBColor(0xd4,0xaf,0x37)},
    }
    theme = GENRE_THEMES.get(data['genre'], {"bg": RGBColor(0x0a,0x0a,0x12), "accent": RGBColor(0x70,0x00,0xff), "title": RGBColor(0xff,0x00,0x3c)})

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    def set_bg(slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_bar(slide, top, height, color):
        bar = slide.shapes.add_shape(1, Inches(0), Inches(top), Inches(13.33), Emu(int(height * 914400)))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()

    def add_text(slide, text, left, top, width, height, size, color, bold=False, italic=False):
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic

    # Title slide
    slide = prs.slides.add_slide(blank)
    set_bg(slide, theme["bg"])
    add_bar(slide, 3.4, 0.06, theme["accent"])
    add_text(slide, f"✦ {data['genre'].upper()} CHRONICLE ✦", 0.5, 0.8, 12, 1.8, 42, theme["title"], bold=True)
    add_text(slide, data['profile']['bio'], 0.5, 2.6, 12, 1.5, 18, RGBColor(0xdd,0xdd,0xdd), italic=True)
    add_text(slide, f"by  {data['profile']['name']}", 0.5, 5.8, 12, 0.9, 20, theme["accent"], bold=True)

    # Story slides — one slide per paragraph
    slide_num = 1
    for msg in data['messages']:
        if msg['role'] == 'user':
            label       = f"◈  {data['profile']['name']}"
            label_color = theme["accent"]
            text_color  = RGBColor(0xff,0xff,0xff)
            is_italic   = False
        else:
            label       = "◈  Narrator"
            label_color = theme["title"]
            text_color  = RGBColor(0xdd,0xdd,0xdd)
            is_italic   = True

        paragraphs = [p.strip() for p in msg['content'].split('\n') if p.strip()]
        if not paragraphs:
            paragraphs = [msg['content']]

        for para in paragraphs:
            slide = prs.slides.add_slide(blank)
            set_bg(slide, theme["bg"])
            add_bar(slide, 0, 0.18, theme["accent"])
            add_text(slide, label, 0.5, 0.28, 12, 0.7, 22, label_color, bold=True)
            add_text(slide, para, 0.5, 1.1, 12.3, 5.9, 19, text_color, italic=is_italic)
            add_text(slide, str(slide_num), 12.6, 7.1, 0.7, 0.38, 12, theme["accent"])
            slide_num += 1

    bio = BytesIO()
    prs.save(bio)
    return bio.getvalue()

# --- UI STYLING ---


def apply_styles(genre):
    image_url = BACKGROUND_IMAGES.get(genre, BACKGROUND_IMAGES["Default"])
    st.markdown(f'''
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Syncopate:wght@400;700&family=Space+Grotesk:wght@300;400;700&display=swap');

    .stApp {{ background: transparent !important; font-family: 'Space Grotesk', sans-serif; color: #FFFFFF; }}
    
    .stApp::before {{
        content: ""; position: fixed; top: 0; left: 0; width: 100vw; height: 100vh;
        background-image: linear-gradient(rgba(0, 0, 0, 0.7), rgba(0, 0, 0, 0.8)), url("{image_url}");
        background-size: cover; 
        background-position: center top;
        background-repeat: no-repeat;
        z-index: -1; 
        transition: background-image 1.5s ease-in-out;
    }}

    blockquote{{
        background-color: rgba(112, 0, 255, 0.1) !important;
        border-left: 5px solid #FF003C !important;
        padding: 10px 20px !important;
        border-radius: 5px;
        margin: 10px 0px !important;
        color: #f0f0f0 !important;
    }}
    [data-testid="stSidebar"] {{ background-color: rgba(5, 5, 8, 0.85) !important; backdrop-filter: blur(20px); border-right: 1px solid rgba(112, 0, 255, 0.3) !important; }}
    
    .stButton>button {{
        width: 100%; background: linear-gradient(90deg, #7000ff, #FF003C) !important; color: #FFFFFF !important; 
        font-family: 'Syncopate', sans-serif !important; border-radius: 8px !important; transition: all 0.4s ease !important;
    }}
    .stButton>button:hover {{ box-shadow: 0px 0px 30px rgba(112, 0, 255, 0.5); transform: scale(1.02); }}

    [data-testid="stVerticalBlockBorderWrapper"] {{
        background-color: rgba(255, 255, 255, 0.03) !important;
        backdrop-filter: blur(12px);
        border: 1px solid rgba(112, 0, 255, 0.3) !important;
        border-radius: 15px !important;
    }}

    [data-testid="stBottom"] {{ display: none !important; }}
    [data-testid="stHeader"] {{ background: transparent !important; }}
    footer {{ visibility: hidden; }}
    </style>
    ''', unsafe_allow_html=True)


# --- SESSION INITIALIZATION ---
if "logged_in" not in st.session_state:
    st.session_state.logged_in    = False
    st.session_state.current_user = None

# ── GATE: show login page until authenticated ──
if not st.session_state.logged_in:
    col_l, col_c, col_r = st.columns([1, 1.6, 1])
    with col_c:
        show_auth_page()
    st.stop()

if "conversations" not in st.session_state:
    st.session_state.conversations = load_data()
if "preferred_genre" not in st.session_state:
    st.session_state.preferred_genre = "Default"
if "current_chat" not in st.session_state:
    st.session_state.current_chat = None
if "app_mode" not in st.session_state:
    st.session_state.app_mode = "Creative"
if "latest_generation" not in st.session_state:
    st.session_state.latest_generation = None
if "is_loading" not in st.session_state:
    st.session_state.is_loading = False
# ==========================================
# --- NEURAL PROCESSING (REFINED) ---
# ==========================================
if st.session_state.is_loading and st.session_state.current_chat:
    active_id = st.session_state.current_chat
    chat_obj = st.session_state.conversations[active_id]

    with st.spinner("The Narrator is consulting the fates..."):
        system_content = (
            f"You are a master {chat_obj['genre']} writer. Story Synopsis: {chat_obj['profile']['bio']}. "
            "STRICT RULES: You ONLY respond to story-writing requests. "
            "If the user asks about politics, real-world current events, controversial figures, "
            "harmful content, or anything unrelated to this fiction story, "
            "you MUST refuse by saying: 'RESTRICTED: I am a Narrator of tales, not a commentator on worldly affairs.'"
        )

        # Build messages including the new user input already in chat_obj["messages"]
        llama_messages = [{"role": "system", "content": system_content}]
        for msg in chat_obj["messages"][-6:]:
            llama_messages.append(
                {"role": msg["role"], "content": msg["content"]})

        try:
            response = ollama.chat(model="llama3", messages=llama_messages)
            output = response["message"]["content"]

            # --- IMPROVED SAFETY LOGIC ---
            if "RESTRICTED" in output or "politics" in output.lower():
                # 1. Alert the user
                st.toast("The Narrator refuses this path.", icon="⚠️")
                st.session_state.latest_generation = "⚠️ This thread was rejected. Please steer back to the story."

                # 2. OPTIONAL: Remove the last 'harmful' user message so it doesn't pollute the context
                if chat_obj["messages"][-1]["role"] == "user":
                    chat_obj["messages"].pop()
            else:
                st.session_state.latest_generation = output

        except Exception as e:
            st.error(f"Engine Error: {e}. Ensure Ollama is running.")

    st.session_state.is_loading = False
    st.rerun()

# --- SIDEBAR ---
with st.sidebar:
    if st.button("🏠 GO TO HOME", use_container_width=True):
        save_data()
        st.session_state.current_chat = None
        st.session_state.preferred_genre = "Default"
        st.rerun()

    st.divider()
    st.markdown("### 🖋️ **Creator**")
    u_name = st.text_input("Alias", value="Co-Writer")
    u_bio = st.text_area("Story's Synopsis",
                         value="A journey through the unknown limits.")
    st.session_state.app_mode = st.selectbox(
        "💡 AI MODE", MODES, index=MODES.index(st.session_state.app_mode))

    st.divider()
    selected_genre = st.selectbox("GENRE FOCUS", GENRES)

    if st.button("Ignite Creativity", use_container_width=True):
        st.session_state.preferred_genre = selected_genre
        new_id = str(uuid.uuid4())
        st.session_state.conversations[new_id] = {
            "genre": selected_genre,
            "mode": st.session_state.app_mode,
            "profile": {"name": u_name, "bio": u_bio},
            "messages": []
        }
        st.session_state.current_chat = new_id
        save_data()
        st.rerun()

    st.divider()
    for cid in list(st.session_state.conversations.keys()):
        data = st.session_state.conversations[cid]
        c1, c2 = st.columns([4, 2])
        with c1:
            if st.button(f"Draft: {data['genre']}", key=f"b_{cid}", use_container_width=True):
                st.session_state.current_chat = cid
                st.session_state.preferred_genre = data['genre']
                st.rerun()
        with c2:
            with st.popover("⚙️"):
                if st.button("⛔ Delete", key=f"d_{cid}", use_container_width=True):
                    del st.session_state.conversations[cid]
                    if st.session_state.current_chat == cid:
                        st.session_state.current_chat = None
                    save_data()
                    st.rerun()
                st.download_button(label="📄 Export", data=generate_docx(
                    data), file_name=f"Manuscript.docx", key=f"w_{cid}")

# --- MAIN CONTENT ---
if st.session_state.current_chat and st.session_state.current_chat in st.session_state.conversations:
    active_id = st.session_state.current_chat
    chat_obj = st.session_state.conversations[active_id]
    apply_styles(chat_obj["genre"])

    st.title(f"Narrative Flow ~ {chat_obj['genre']}")

    col_left, col_right = st.columns([1.3, 1], gap="large")

    # --- LEFT SIDE: MANUSCRIPT ---
    with col_left:
        st.markdown("### 📜 Manuscript")
        chat_display = st.container(height=350, border=True)
        with chat_display:
            if not chat_obj["messages"]:
                st.info("The parchment is blank. Begin your story below...")
            for m in chat_obj["messages"]:
                if m["role"] == "user":
                    st.markdown(
                        f"**{chat_obj['profile']['name']}:** {m['content']}")
                else:
                    st.markdown(f"> **Narrator:** {m['content']}")

        # Input Form
        with st.form("input_form", clear_on_submit=True):
            user_input = st.text_input(
                "Next line...", placeholder="Type here...")
            submit_weave = st.form_submit_button(
                "WEAVE NARRATIVE →", use_container_width=True)

            if submit_weave and user_input:
                chat_obj["messages"].append(
                    {"role": "user", "content": user_input})
                st.session_state.is_loading = True
                save_data()
                st.rerun()

        # Export Buttons
        exp_col1, exp_col2, exp_col3 = st.columns(3)
        with exp_col1:
            final_story = "\n\n".join(
                [f"{m['role'].upper()}: {m['content']}" for m in chat_obj["messages"]])
            st.download_button(
                label="📄 TXT",
                data=final_story,
                file_name=f"Story_{chat_obj['genre']}.txt",
                mime="text/plain",
                use_container_width=True
            )
        with exp_col2:
            st.download_button(
                label="📝 DOCX",
                data=generate_docx(chat_obj),
                file_name=f"Story_{chat_obj['genre']}.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                use_container_width=True
            )
        with exp_col3:
            st.download_button(
                label="🎞️ PPT",
                data=generate_pptx(chat_obj),
                file_name=f"Story_{chat_obj['genre']}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True
            )

    # --- RIGHT SIDE: AI WORKSHOP ---
    with col_right:
        st.markdown("### ✨ AI Workshop")

        with st.container(height=400, border=True):
            st.markdown("**Generated Suggestion:**")
            if st.session_state.latest_generation:
                st.write(st.session_state.latest_generation)
            else:
                st.info(
                    "The AI's next suggestion will appear here after you write or click Recommend.")

        st.divider()

        btn_col1, btn_col2 = st.columns(2)
        with btn_col1:
            with st.container(border=True):
                st.markdown("📝 **Finalize**")
                if st.button("Save Story", use_container_width=True, key="save_btn"):
                    if st.session_state.latest_generation and "⚠️" not in st.session_state.latest_generation:
                        chat_obj["messages"].append(
                            {"role": "assistant", "content": st.session_state.latest_generation})
                        st.session_state.latest_generation = None
                        save_data()
                        st.rerun()
                    else:
                        st.warning("Cannot save restricted or empty content.")

        with btn_col2:
            with st.container(border=True):
                st.markdown("🔄 **Revision**")
                if st.button("Recommend", use_container_width=True, key="rec_btn"):
                    st.session_state.is_loading = True
                    st.rerun()

else:
    # --- DEFAULT HOME PAGE ---
    apply_styles("Default")
    st.title(" 🌌 NARRATIVE SANCTUARY ")
    st.markdown(
        "#### *Select a draft from the sidebar or explore these creative references.*")

    st.divider()

    ref_col1, ref_col2, ref_col3 = st.columns(3)

    with ref_col1:
        with st.container(border=True):
            st.markdown("### 📚 Word Engines")
            st.markdown(
                "- [**OneLook Thesaurus**](https://www.onelook.com/thesaurus/)")
            st.markdown("- [**Etymonline**](https://www.etymonline.com/)")
            st.markdown(
                "- [**Power Thesaurus**](https://www.powerthesaurus.org/)")

    with ref_col2:
        with st.container(border=True):
            st.markdown("### 🗺️ World Building")
            st.markdown("- [**World Anvil**](https://www.worldanvil.com/)")
            st.markdown(
                "- [**Azgaar's Map**](https://azgaar.github.io/Fantasy-Map-Generator/)")
            st.markdown("- [**Donjon**](https://donjon.bin.sh/)")

    with ref_col3:
        with st.container(border=True):
            st.markdown("### ✍️ Writing Help")
            st.markdown(
                "- [**Helping Writers**](https://www.helpingwritersbecomeauthors.com/)")
            st.markdown("- [**TV Tropes**](https://tvtropes.org/)")
            st.markdown(
                "- [**Reedsy Learning**](https://blog.reedsy.com/learning/)")

    st.divider()
    st.info("💡 **Pro-Tip:** Use the sidebar to set your **Alias** and **Synopsis** before clicking 'Ignite Creativity' for the best AI-guided results.")
