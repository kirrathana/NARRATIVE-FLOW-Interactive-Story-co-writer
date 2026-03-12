import io as _io
import base64
import streamlit as st
import streamlit.components.v1 as components
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from ollama_helper import check_ollama_connection, generate_story_stream, get_model_name, check_guardrails
from db import (
    init_db, create_user, verify_user, save_story,
    get_stories, get_story_by_id, delete_story, get_user_info, format_date,
)

st.set_page_config(
    page_title="Story Waver",
    page_icon="📖",
    layout="wide",
    initial_sidebar_state="expanded",
)

init_db()

# ── CSS ───────────────────────────────────────────────────────────────────────
st.markdown("""
<style>
/* ── Global ── */
.stApp { background-color: #0f0f1a; color: #e0e0e0; }
[data-testid="stSidebar"] { background-color: #13131f; border-right: 1px solid #1e1e35; }
[data-testid="stSidebar"] * { color: #c8c8d8 !important; }

/* ── Title ── */
.main-title {
    text-align: center; font-size: 2.6rem; font-weight: 900; letter-spacing: 8px;
    background: linear-gradient(135deg, #00d4aa 0%, #0099cc 50%, #6633cc 100%);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text; margin-bottom: 0; padding-top: 10px;
}
.main-subtitle {
    text-align: center; color: #555577; font-size: 0.9rem;
    margin-top: 2px; margin-bottom: 10px; letter-spacing: 1px;
}

/* ── Chat messages ── */
[data-testid="stChatMessage"] {
    background: transparent !important;
    border: none !important;
    padding: 4px 0 !important;
}
[data-testid="stChatMessageContent"] {
    background: linear-gradient(145deg, #13131f, #1a1a2e) !important;
    border: 1px solid #252545 !important;
    border-radius: 12px !important;
    padding: 18px 22px !important;
    color: #dde0f0 !important;
    line-height: 1.85 !important;
    font-size: 1.0rem !important;
}
/* User bubble — slightly different tint */
[data-testid="stChatMessage"]:has([data-testid="stChatMessageAvatarUser"]) [data-testid="stChatMessageContent"] {
    background: linear-gradient(145deg, #161628, #1e1e38) !important;
    border-color: #2a2a55 !important;
}

/* ── Chat input bar ── */
[data-testid="stChatInput"] {
    background: #13131f !important;
    border: 1px solid #2a2a4a !important;
    border-radius: 14px !important;
}
[data-testid="stChatInput"] textarea {
    background: transparent !important;
    color: #dde0f0 !important;
    font-size: 0.97rem !important;
}
[data-testid="stChatInput"] button {
    background: linear-gradient(135deg, #00d4aa, #0099cc) !important;
    border-radius: 8px !important;
}
/* Bottom padding so last message isn't hidden behind input */
.main > div { padding-bottom: 90px !important; }

/* ── Genre / Tone tags ── */
.msg-tag {
    display: inline-block; font-size: 0.66rem; font-weight: 700;
    letter-spacing: 1px; text-transform: uppercase;
    padding: 2px 8px; border-radius: 20px; margin-right: 5px; margin-bottom: 8px;
}
.tag-genre { background: #0d2a22; color: #00d4aa !important; }
.tag-tone  { background: #1a1a2e; color: #8888cc !important; border: 1px solid #333360; }
.tag-mode  { background: #1a1a2e; color: #cc8844 !important; border: 1px solid #443322; }

/* ── Sidebar section header ── */
.sidebar-section-header {
    background: linear-gradient(90deg, #1a1a2e, #16213e);
    border-left: 3px solid #00d4aa; padding: 6px 10px; border-radius: 4px;
    font-size: 0.73rem; font-weight: 700; letter-spacing: 1.5px;
    text-transform: uppercase; color: #00d4aa !important;
    margin-bottom: 8px; margin-top: 12px;
}

/* ── Sidebar history card ── */
.history-card {
    background: #1a1a2e; border: 1px solid #252545; border-radius: 10px;
    padding: 10px 12px; margin-bottom: 8px;
    transition: border-color 0.2s;
}
.history-card:hover { border-color: #00d4aa; }
.history-prompt-label { color: #00d4aa !important; font-size: 0.68rem; font-weight: 600; letter-spacing: 0.04em; text-transform: uppercase; margin-top: 6px; margin-bottom: 2px; }
.history-title { color: #ccccee !important; font-size: 0.82rem; line-height: 1.4; margin-bottom: 3px; }
.history-date  { color: #44445a !important; font-size: 0.7rem; }
.reading-time  { color: #44445a !important; font-size: 0.68rem; margin-left: 6px; }

/* ── Buttons ── */
.stButton > button {
    background: linear-gradient(135deg, #1a1a2e, #16213e);
    color: #c0c8e0 !important; border: 1px solid #2a2a4a;
    border-radius: 8px; font-size: 0.82rem; transition: all 0.2s; width: 100%;
}
.stButton > button:hover { border-color: #00d4aa; color: #00d4aa !important; }
.stButton > button[kind="primary"] {
    background: linear-gradient(135deg, #00d4aa, #0099cc);
    color: #0a0a14 !important; font-weight: 700; border: none;
}
.stButton > button[kind="primary"]:hover {
    background: linear-gradient(135deg, #00e8bb, #00aadd);
    box-shadow: 0 4px 20px rgba(0,212,170,0.4);
}

/* ── User badge ── */
.user-badge {
    background: #0d2a22; border: 1px solid #1a5540; border-radius: 20px;
    padding: 5px 14px; font-size: 0.82rem; color: #00d4aa !important; font-weight: 600;
}

/* ── Status ── */
.status-online  { color: #00d4aa; font-weight: 700; font-size: 0.76rem; }
.status-offline { color: #ff4466; font-weight: 700; font-size: 0.76rem; }

/* ── Input fields ── */
.stTextArea textarea, .stTextInput input {
    background-color: #1a1a2e !important; color: #dde0f0 !important;
    border: 1px solid #2a2a4a !important; border-radius: 8px !important;
}

/* ── History page ── */
.story-view-header {
    background: linear-gradient(90deg, #0d2a22, #0d1a2a);
    border: 1px solid #1a4a38; border-radius: 10px;
    padding: 14px 18px; margin-bottom: 16px;
}
.story-view-prompt { color: #c8d8f0 !important; font-size: 0.97rem; margin-top: 8px; line-height: 1.5; }
.story-view-prompt span { color: #7ab8e8 !important; font-weight: 600; margin-right: 4px; }
.story-view-meta { color: #6688aa !important; font-size: 0.78rem; margin-top: 6px; }
.story-box {
    background: linear-gradient(145deg, #13131f, #1a1a2e);
    border: 1px solid #252545; border-radius: 12px;
    padding: 24px 28px; line-height: 1.85; font-size: 1rem;
    color: #dde0f0; white-space: pre-wrap;
}

hr { border-color: #252545 !important; }
[data-testid="stExpander"] { background: #13131f; border: 1px solid #252545 !important; border-radius: 8px; }

/* ── Stats Row ── */
.stats-row { display: flex; gap: 12px; margin: 0 0 24px; }
.stat-card {
    flex: 1; background: #13131f; border: 1px solid #252545;
    border-radius: 14px; padding: 18px 12px; text-align: center;
    transition: border-color 0.3s, transform 0.2s;
}
.stat-card:hover { border-color: #00d4aa; transform: translateY(-2px); }
.stat-icon { font-size: 1.4rem; display: block; margin-bottom: 8px; }
.stat-value {
    font-size: 1.9rem; font-weight: 800; display: block;
    background: linear-gradient(135deg, #00d4aa, #0099cc);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text; line-height: 1.2; margin-bottom: 5px;
}
.stat-label { color: #333355 !important; font-size: 0.63rem; font-weight: 700; letter-spacing: 1.8px; text-transform: uppercase; display: block; }

/* ── Story Progress Bar ── */
.progress-wrap {
    background: #13131f; border: 1px solid #252545; border-radius: 14px;
    padding: 18px 24px; margin-bottom: 24px;
}
.progress-title { color: #444466 !important; font-size: 0.68rem; font-weight: 700; letter-spacing: 2px; text-transform: uppercase; margin-bottom: 14px; }
.progress-phases { display: flex; align-items: center; }
.phase-node {
    display: flex; flex-direction: column; align-items: center; gap: 6px; flex: 0 0 auto;
}
.phase-dot {
    width: 12px; height: 12px; border-radius: 50%;
    background: #1e1e35; border: 2px solid #252545;
    transition: all 0.3s;
}
.phase-dot.active { background: #00d4aa; border-color: #00d4aa; box-shadow: 0 0 10px #00d4aa66; }
.phase-dot.done   { background: #005544; border-color: #007755; }
.phase-label { font-size: 0.68rem; color: #333355 !important; font-weight: 600; letter-spacing: 0.5px; }
.phase-label.active { color: #00d4aa !important; }
.phase-label.done   { color: #007755 !important; }
.phase-line { flex: 1; height: 2px; background: #1e1e35; margin: 0 6px; margin-bottom: 18px; border-radius: 2px; }
.phase-line.done { background: linear-gradient(90deg, #005544, #00d4aa44); }
.progress-words { color: #444466 !important; font-size: 0.72rem; margin-top: 12px; font-weight: 500; }

/* ── Genre Background Overlay ── */
.genre-overlay {
    position: fixed; inset: 0; pointer-events: none; z-index: 1; overflow: hidden;
}
.gb-elem {
    position: absolute; bottom: -60px; opacity: 0; user-select: none;
    animation-iteration-count: infinite; animation-timing-function: linear;
}
@keyframes gbFloatUp {
    0%   { transform: translateY(0) rotate(0deg);       opacity: 0; }
    8%   { opacity: var(--gb-opacity, 0.15); }
    92%  { opacity: var(--gb-opacity, 0.15); }
    100% { transform: translateY(-110vh) rotate(25deg);  opacity: 0; }
}
@keyframes gbDrift {
    0%   { transform: translate(0,0) rotate(0deg);           opacity: 0; }
    8%   { opacity: var(--gb-opacity, 0.15); }
    40%  { transform: translate(40px,-45vh) rotate(180deg);  opacity: var(--gb-opacity, 0.15); }
    60%  { transform: translate(-30px,-65vh) rotate(270deg); opacity: var(--gb-opacity, 0.15); }
    92%  { transform: translate(10px,-100vh) rotate(340deg); opacity: var(--gb-opacity, 0.15); }
    100% { transform: translate(10px,-110vh) rotate(360deg); opacity: 0; }
}
@keyframes gbTwinkle {
    0%,100% { opacity: 0;                      transform: scale(0.5) translateY(0);    }
    30%,70% { opacity: var(--gb-opacity, 0.2); transform: scale(1.1) translateY(-10px); }
}
@keyframes gbBounce {
    0%   { transform: translateY(0) scale(1) rotate(0deg);    opacity: 0; }
    5%   { opacity: var(--gb-opacity, 0.15); }
    20%  { transform: translateY(-20vh) scale(1.1) rotate(10deg); }
    50%  { transform: translateY(-55vh) scale(0.9) rotate(-10deg); }
    80%  { transform: translateY(-85vh) scale(1.05) rotate(5deg); }
    95%  { opacity: var(--gb-opacity, 0.15); }
    100% { transform: translateY(-110vh) scale(0.8) rotate(0deg); opacity: 0; }
}
@keyframes gbPulse {
    0%,100% { opacity: 0.03; transform: scale(0.9); }
    50%     { opacity: var(--gb-opacity, 0.18); transform: scale(1.15) translateY(-5px); }
}
@keyframes gbSway {
    0%   { transform: translateY(0) rotate(-8deg);     opacity: 0; }
    8%   { opacity: var(--gb-opacity, 0.15); }
    25%  { transform: translateY(-28vh) rotate(8deg);  opacity: var(--gb-opacity, 0.15); }
    50%  { transform: translateY(-55vh) rotate(-8deg); opacity: var(--gb-opacity, 0.15); }
    75%  { transform: translateY(-80vh) rotate(8deg);  opacity: var(--gb-opacity, 0.15); }
    92%  { opacity: var(--gb-opacity, 0.15); }
    100% { transform: translateY(-110vh) rotate(-8deg); opacity: 0; }
}
</style>
""", unsafe_allow_html=True)


# ── Session State ─────────────────────────────────────────────────────────────
for key, default in {
    "logged_in": False,
    "username": "",
    "messages": [],           # chat turns: {role, content, genre, tone, writing_mode}
    "view_story_id": None,
    "page": "main",
}.items():
    if key not in st.session_state:
        st.session_state[key] = default


# ── User stats helper ─────────────────────────────────────────────────────────
def _get_user_stats(username: str):
    stories = get_stories(username)
    total_stories = len(stories)
    total_words   = sum(len(s["story"].split()) for s in stories)
    genre_counts  = {}
    for s in stories:
        if s["genre"] != "Mixed":
            genre_counts[s["genre"]] = genre_counts.get(s["genre"], 0) + 1
    fav_genre = max(genre_counts, key=genre_counts.get) if genre_counts else "—"
    return total_stories, total_words, fav_genre


# ── Export helpers ─────────────────────────────────────────────────────────────
_GENRE_ACCENT = {
    "Mixed":     (  0, 212, 170),
    "Romantic":  (220,  60, 110),
    "Horror":    ( 20, 200,  50),
    "Fantasy":   (130,  80, 220),
    "Mystery":   ( 60,  90, 200),
    "Humor":     (220, 170,  20),
    "Suspense":  (210,  30,  30),
    "Sci-Fi":    ( 20, 130, 230),
    "Adventure": ( 50, 170,  50),
}
_GENRE_BG_RGB = {
    "Mixed":     (10,  15,  25),
    "Romantic":  (45,  10,  30),
    "Horror":    (10,  10,  10),
    "Fantasy":   (16,   8,  48),
    "Mystery":   ( 9,  13,  26),
    "Humor":     (20,  18,   4),
    "Suspense":  (31,   3,   3),
    "Sci-Fi":    ( 2,  10,  31),
    "Adventure": ( 6,  14,   2),
}

def _safe_text(text: str) -> str:
    """Transliterate common Unicode to latin-1 for FPDF built-in fonts."""
    subs = {
        "\u2018": "'", "\u2019": "'", "\u201c": '"', "\u201d": '"',
        "\u2013": "-", "\u2014": "--", "\u2026": "...", "\u00a0": " ",
        "\u2022": "*", "\u2032": "'",
    }
    for k, v in subs.items():
        text = text.replace(k, v)
    return text.encode("latin-1", errors="replace").decode("latin-1")


class _StoryPDF(FPDF):
    def __init__(self, accent_rgb, genre):
        super().__init__()
        self._ar, self._ag, self._ab = accent_rgb
        self._genre = genre

    def footer(self):
        self.set_y(-12)
        self.set_draw_color(self._ar, self._ag, self._ab)
        self.line(10, self.get_y(), 200, self.get_y())
        self.set_y(-10)
        self.set_font("Helvetica", "", 7)
        self.set_text_color(150, 150, 160)
        self.cell(95, 5, "Generated by Story Waver")
        self.cell(95, 5, f"Page {self.page_no()}", align="R")


@st.cache_data(show_spinner=False)
def make_pdf(sid, title, genre, tone, mode, prompt, story_text, created_at) -> bytes:
    accent = _GENRE_ACCENT.get(genre, (0, 212, 170))
    r, g, b = accent

    pdf = _StoryPDF(accent, genre)
    pdf.set_auto_page_break(auto=True, margin=18)
    pdf.add_page()

    # Header bar
    pdf.set_fill_color(r, g, b)
    pdf.rect(0, 0, 210, 14, "F")
    pdf.set_xy(8, 3)
    pdf.set_font("Helvetica", "B", 8)
    pdf.set_text_color(255, 255, 255)
    pdf.cell(0, 8, _safe_text(f"STORY WAVER  |  {genre.upper()}  |  {tone.upper()}  |  {mode.upper()}"))

    # Title
    pdf.set_xy(10, 20)
    pdf.set_font("Helvetica", "B", 20)
    pdf.set_text_color(30, 30, 50)
    pdf.multi_cell(190, 10, _safe_text(title))

    # Prompt
    pdf.set_x(10)
    pdf.set_font("Helvetica", "I", 9)
    pdf.set_text_color(110, 110, 140)
    pdf.multi_cell(190, 5, _safe_text(f'Prompt: "{prompt}"'))

    # Date
    pdf.set_x(10)
    pdf.set_font("Helvetica", "", 8)
    pdf.set_text_color(150, 150, 160)
    pdf.cell(0, 6, _safe_text(format_date(created_at)))

    # Divider
    y = pdf.get_y() + 4
    pdf.set_draw_color(r, g, b)
    pdf.set_line_width(0.6)
    pdf.line(10, y, 200, y)

    # Story body
    pdf.set_xy(10, y + 6)
    pdf.set_font("Helvetica", "", 11)
    pdf.set_text_color(30, 30, 50)
    pdf.multi_cell(190, 6.5, _safe_text(story_text))

    return bytes(pdf.output())


@st.cache_data(show_spinner=False)
def make_pptx(sid, title, genre, tone, mode, prompt, story_text, created_at) -> bytes:
    ar, ag, ab = _GENRE_ACCENT.get(genre, (0, 212, 170))
    br, bg_, bb = _GENRE_BG_RGB.get(genre, (15, 15, 26))
    accent = RGBColor(ar, ag, ab)
    bg_col = RGBColor(br, bg_, bb)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    W, H  = prs.slide_width, prs.slide_height

    def _bg(slide):
        s = slide.shapes.add_shape(1, 0, 0, W, H)
        s.fill.solid(); s.fill.fore_color.rgb = bg_col; s.line.fill.background()
        bar = slide.shapes.add_shape(1, 0, 0, W, Pt(6))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()

    def _txt(slide, l, t, w, h, text, size, bold=False, italic=False,
             color=None, align=PP_ALIGN.LEFT, wrap=True):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame; tf.word_wrap = wrap
        p  = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = text
        run.font.size = Pt(size); run.font.bold = bold; run.font.italic = italic
        run.font.color.rgb = color or RGBColor(210, 214, 230)

    # ── Title slide ──
    sl = prs.slides.add_slide(blank)
    _bg(sl)
    _txt(sl, Inches(0.5), Inches(0.2),  Inches(5),  Inches(0.4),
         "STORY WAVER", 9, bold=True, color=accent)
    _txt(sl, Inches(0.5), Inches(0.65), Inches(12), Inches(0.4),
         f"{genre}  ·  {tone}  ·  {mode}", 10, color=RGBColor(160, 165, 185))
    _txt(sl, Inches(0.5), Inches(1.6),  Inches(12), Inches(2.2),
         title, 34, bold=True, wrap=True)
    line = sl.shapes.add_shape(1, Inches(0.5), Inches(4.15), Inches(2.5), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = accent; line.line.fill.background()
    _txt(sl, Inches(0.5), Inches(4.45), Inches(12), Inches(0.7),
         f'"{prompt}"', 11, italic=True, color=RGBColor(130, 150, 140))
    _txt(sl, Inches(0.5), Inches(5.2),  Inches(6),  Inches(0.4),
         format_date(created_at), 9, color=RGBColor(90, 90, 110))

    # ── Story content slides ──
    paras = [p for p in story_text.split("\n") if p.strip()]
    chunks, cur = [], ""
    for para in paras:
        if len(cur) + len(para) + 2 < 650:
            cur += ("\n\n" if cur else "") + para
        else:
            if cur: chunks.append(cur)
            cur = para
    if cur:
        chunks.append(cur)

    total = len(chunks) + 1
    for idx, chunk in enumerate(chunks):
        sl = prs.slides.add_slide(blank)
        _bg(sl)
        _txt(sl, Inches(10.8), Inches(0.12), Inches(2.3), Inches(0.3),
             genre.upper(), 8, bold=True, color=accent,
             align=PP_ALIGN.RIGHT, wrap=False)
        _txt(sl, Inches(0.6), Inches(0.55), Inches(12.1), Inches(6.55),
             chunk, 14, wrap=True)
        _txt(sl, Inches(11.3), Inches(7.1), Inches(1.8), Inches(0.3),
             f"{idx + 2} / {total}", 8,
             color=RGBColor(80, 80, 100), align=PP_ALIGN.RIGHT, wrap=False)

    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Genre Background Themes ────────────────────────────────────────────────────
def apply_genre_background(genre: str) -> None:
    """Inject floating emoji overlay + background gradient for the selected genre."""
    CONFIGS = {
        "Mixed": {
            "bg":      "radial-gradient(ellipse at 50% 60%, #0a0f19 0%, #080c14 55%, #050810 100%)",
            "emojis":  ["📖", "✨", "🌙", "⚔️", "🔍", "🚀", "❤️", "👻", "🎭", "🌟"],
            "anim":    "gbFloatUp", "opacity": 0.18, "count": 16,
        },
        "Romantic": {
            "bg":      "radial-gradient(ellipse at 50% 85%, #2d0a1e 0%, #1a060f 55%, #0d0308 100%)",
            "emojis":  ["❤️", "💕", "💗", "🌹", "💝", "💖", "💘", "🌸", "💞", "🥀"],
            "anim":    "gbFloatUp", "opacity": 0.22, "count": 18,
        },
        "Horror": {
            "bg":      "radial-gradient(ellipse at 50% 50%, #0c0c0c 0%, #050505 55%, #000000 100%)",
            "emojis":  ["👻", "💀", "🕸️", "🦇", "🕷️", "⚰️", "🌑", "🪦", "😱", "🩸"],
            "anim":    "gbDrift", "opacity": 0.22, "count": 14,
        },
        "Fantasy": {
            "bg":      "radial-gradient(ellipse at 50% 50%, #100830 0%, #080520 55%, #040212 100%)",
            "emojis":  ["✨", "🌟", "⭐", "🦋", "🌙", "🔮", "🧚", "🌈", "🦄", "🧝", "🪄"],
            "anim":    "gbTwinkle", "opacity": 0.28, "count": 22,
        },
        "Mystery": {
            "bg":      "radial-gradient(ellipse at 50% 50%, #090d1a 0%, #060a12 55%, #030508 100%)",
            "emojis":  ["🔍", "🕵️", "🗝️", "📋", "🔎", "❓", "🎩", "🕶️", "📜", "🔦"],
            "anim":    "gbFloatUp", "opacity": 0.18, "count": 12,
        },
        "Humor": {
            "bg":      "radial-gradient(ellipse at 50% 50%, #0d0c02 0%, #0a0902 55%, #070602 100%)",
            "emojis":  ["😂", "🤣", "😆", "😄", "🎭", "🎪", "🃏", "🎈", "😝", "🤡", "🥳", "😜"],
            "anim":    "gbBounce", "opacity": 0.24, "count": 18,
        },
        "Suspense": {
            "bg":      "radial-gradient(ellipse at 50% 50%, #1f0303 0%, #120202 55%, #080101 100%)",
            "emojis":  ["🚨", "🩸", "⚠️", "🚔", "👁️", "🔪", "🕵️", "😰", "⏰", "🔫"],
            "anim":    "gbPulse", "opacity": 0.22, "count": 12,
        },
        "Sci-Fi": {
            "bg":      "radial-gradient(ellipse at 50% 50%, #020a1f 0%, #010610 55%, #000308 100%)",
            "emojis":  ["🚀", "⭐", "🌟", "🛸", "🌌", "👾", "🤖", "🌠", "🛰️", "🔭", "👽"],
            "anim":    "gbTwinkle", "opacity": 0.24, "count": 20,
        },
        "Adventure": {
            "bg":      "radial-gradient(ellipse at 50% 85%, #060e02 0%, #040a01 55%, #020601 100%)",
            "emojis":  ["⚔️", "🗺️", "🧭", "🏔️", "🌿", "🦅", "💎", "🏹", "🌋", "⛵", "🦁"],
            "anim":    "gbSway", "opacity": 0.20, "count": 14,
        },
    }

    cfg = CONFIGS.get(genre)
    if not cfg:
        return

    emojis  = cfg["emojis"]
    n       = cfg["count"]
    opacity = cfg["opacity"]
    anim    = cfg["anim"]
    scatter = anim in ("gbTwinkle", "gbPulse")

    spans = ""
    for i in range(n):
        left  = round((i * (100.0 / n) + (i % 3) * 5.5) % 97, 1)
        delay = round((i * 1.73) % 13, 1)
        dur   = 8 + (i % 7) * 2
        size  = 22 + (i % 5) * 10
        em    = emojis[i % len(emojis)]
        pos   = f"top:{(i * 17 + 5) % 90}%;bottom:auto;" if scatter else ""
        spans += (
            f"<span class='gb-elem' style='"
            f"left:{left}%;{pos}"
            f"animation-name:{anim};"
            f"animation-delay:{delay}s;"
            f"animation-duration:{dur}s;"
            f"font-size:{size}px;'>"
            f"{em}</span>"
        )

    st.markdown(
        f"<style>.stApp{{background:{cfg['bg']} !important;}}</style>"
        f"<div class='genre-overlay' style='--gb-opacity:{opacity};'>{spans}</div>",
        unsafe_allow_html=True,
    )


# ══════════════════════════════════════════════════════════════════════════════
#  AUTH PAGE
# ══════════════════════════════════════════════════════════════════════════════
def show_auth_page():
    st.markdown("<div class='main-title'>STORY WAVER</div>", unsafe_allow_html=True)
    st.markdown("<div class='main-subtitle'>Your personal AI story companion ✨</div>", unsafe_allow_html=True)

    # Inject JS so the browser password manager can autofill saved credentials
    components.html("""
    <script>
    (function() {
        function applyAutocomplete() {
            var doc = window.parent.document;
            doc.querySelectorAll('[data-testid="stTextInput"]').forEach(function(container) {
                var label = container.querySelector('label');
                var input = container.querySelector('input');
                if (!label || !input) return;
                var text = label.textContent.toLowerCase().trim();
                if (text === 'username' || text === 'choose username') {
                    input.setAttribute('autocomplete', text === 'choose username' ? 'new-username' : 'username');
                    input.setAttribute('name', text === 'choose username' ? 'new-username' : 'username');
                } else if (text === 'password' || text === 'choose password') {
                    input.setAttribute('autocomplete', text === 'choose password' ? 'new-password' : 'current-password');
                    input.setAttribute('name', text === 'choose password' ? 'new-password' : 'password');
                } else if (text.includes('confirm')) {
                    input.setAttribute('autocomplete', 'new-password');
                    input.setAttribute('name', 'confirm-password');
                }
            });
        }
        applyAutocomplete();
        setInterval(applyAutocomplete, 400);
    })();
    </script>
    """, height=0)

    _, col, _ = st.columns([1, 1.4, 1])
    with col:
        tab_login, tab_register = st.tabs(["Login", "Register"])

        with tab_login:
            st.markdown("<br>", unsafe_allow_html=True)
            with st.form("login_form"):
                uname = st.text_input("Username", key="login_user", placeholder="your username")
                pwd   = st.text_input("Password", type="password", key="login_pwd", placeholder="••••••••")
                st.markdown("<br>", unsafe_allow_html=True)
                submitted = st.form_submit_button("Login", type="primary", use_container_width=True)
            if submitted:
                if uname and pwd:
                    ok, msg = verify_user(uname, pwd)
                    if ok:
                        st.session_state.logged_in = True
                        st.session_state.username = uname.strip().lower()
                        st.session_state.messages = []
                        st.rerun()
                    else:
                        st.error(msg)
                else:
                    st.warning("Please enter username and password.")

        with tab_register:
            st.markdown("<br>", unsafe_allow_html=True)
            with st.form("register_form"):
                new_user = st.text_input("Choose Username", key="reg_user",  placeholder="min 3 characters")
                new_pwd  = st.text_input("Choose Password", type="password", key="reg_pwd",  placeholder="min 4 characters")
                new_pwd2 = st.text_input("Confirm Password", type="password", key="reg_pwd2", placeholder="repeat password")
                st.markdown("<br>", unsafe_allow_html=True)
                reg_submitted = st.form_submit_button("Create Account", type="primary", use_container_width=True)
            if reg_submitted:
                if new_pwd != new_pwd2:
                    st.error("Passwords do not match.")
                elif new_user and new_pwd:
                    ok, msg = create_user(new_user, new_pwd)
                    if ok:
                        st.success(msg + " Please login.")
                    else:
                        st.error(msg)
                else:
                    st.warning("Please fill in all fields.")


# ══════════════════════════════════════════════════════════════════════════════
#  HISTORY PAGE
# ══════════════════════════════════════════════════════════════════════════════
def show_history_page():
    username = st.session_state.username
    stories  = get_stories(username)

    st.markdown("<div class='main-title' style='font-size:2rem;'>MY STORIES</div>", unsafe_allow_html=True)
    st.markdown(
        f"<div class='main-subtitle'>{len(stories)} stories saved for "
        f"<b style='color:#00d4aa'>@{username}</b></div>",
        unsafe_allow_html=True,
    )

    if st.button("← Back to Writer", key="back_from_history"):
        st.session_state.view_story_id = None
        st.session_state.page = "main"
        st.rerun()

    st.markdown("---")

    # ── Single story view ──
    if st.session_state.view_story_id:
        story = get_story_by_id(username, st.session_state.view_story_id)
        if story:
            st.markdown(
                f"<div class='story-view-header'>"
                f"<span class='msg-tag tag-genre'>{story['genre']}</span>"
                f"<span class='msg-tag tag-mode'>{story['writing_mode']}</span>"
                f"<span class='msg-tag tag-tone'>{story['tone']}</span>"
                f"<div style='color:#ccccee;font-size:0.95rem;font-weight:500;margin-top:10px;line-height:1.4;'>{story['title']}</div>"
                f"<div class='story-view-prompt'><span>Prompt:</span>{story['prompt']}</div>"
                f"<div class='story-view-meta'>{format_date(story['created_at'])}</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
            st.markdown(f"<div class='story-box'>{story['story']}</div>", unsafe_allow_html=True)
            c1, c2 = st.columns([3, 2])
            with c1:
                if st.button("← Back to History", key="back_to_list"):
                    st.session_state.view_story_id = None
                    st.rerun()
            with c2:
                if st.button("Load into Writer", type="primary", key="load_story", use_container_width=True):
                    # Inject as a message so it appears in chat
                    st.session_state.messages = [
                        {"role": "user",      "content": story["prompt"]},
                        {"role": "assistant", "content": story["story"],
                         "genre": story["genre"], "writing_mode": story["writing_mode"],
                         "tone": story["tone"]},
                    ]
                    st.session_state.view_story_id = None
                    st.session_state.page = "main"
                    st.rerun()

            # ── Export row ──
            st.markdown(
                "<div style='margin-top:10px;margin-bottom:4px;color:#555577;"
                "font-size:0.75rem;letter-spacing:1.2px;text-transform:uppercase;"
                "font-weight:700;'>⬇ Export Story As</div>",
                unsafe_allow_html=True,
            )
            ex1, ex2, ex3 = st.columns(3)
            sid = story["id"]
            with ex1:
                st.download_button(
                    "📄 PDF",
                    data=make_pdf(sid, story["title"], story["genre"], story["tone"],
                                  story["writing_mode"], story["prompt"],
                                  story["story"], story["created_at"]),
                    file_name=f"story_{sid}.pdf",
                    mime="application/pdf",
                    use_container_width=True,
                    key=f"dl_pdf_{sid}",
                )
            with ex2:
                st.download_button(
                    "📊 PPT",
                    data=make_pptx(sid, story["title"], story["genre"], story["tone"],
                                   story["writing_mode"], story["prompt"],
                                   story["story"], story["created_at"]),
                    file_name=f"story_{sid}.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True,
                    key=f"dl_ppt_{sid}",
                )
            with ex3:
                st.download_button(
                    "📝 TXT",
                    data=story["story"],
                    file_name=f"story_{sid}.txt",
                    mime="text/plain",
                    use_container_width=True,
                    key=f"dl_txt_{sid}",
                )
        return

    # ── Story list ──
    if not stories:
        st.markdown(
            "<div style='text-align:center;color:#333355;padding:60px 0;font-size:1.1rem;'>"
            "📜 No stories yet — go write your first one!</div>",
            unsafe_allow_html=True,
        )
        return

    all_genres   = sorted(set(s["genre"] for s in stories))
    genre_filter = st.selectbox("Filter by Genre", ["All"] + all_genres, key="hist_filter")
    filtered     = stories if genre_filter == "All" else [s for s in stories if s["genre"] == genre_filter]

    st.markdown(f"<div style='color:#555577;font-size:0.8rem;margin-bottom:12px;'>{len(filtered)} stories</div>",
                unsafe_allow_html=True)

    for story in filtered:
        c1, c2 = st.columns([6, 1])
        with c1:
            _prompt_snippet = story['prompt'][:90] + "…" if len(story['prompt']) > 90 else story['prompt']
            st.markdown(
                f"<div class='history-card'>"
                f"<span class='msg-tag tag-genre'>{story['genre']}</span>"
                f"<span class='msg-tag tag-tone'>{story['tone']}</span>"
                f"<div class='history-prompt-label'>Prompt</div>"
                f"<div class='history-title'>{_prompt_snippet}</div>"
                f"<div class='history-date'>{format_date(story['created_at'])}</div>"
                f"</div>",
                unsafe_allow_html=True,
            )
        with c2:
            st.markdown("<div style='padding-top:8px;'>", unsafe_allow_html=True)
            if st.button("View",   key=f"view_{story['id']}", use_container_width=True):
                st.session_state.view_story_id = story["id"]
                st.rerun()
            if st.button("Delete", key=f"del_{story['id']}",  use_container_width=True):
                delete_story(username, story["id"])
                st.rerun()
            st.markdown("</div>", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════════════════
#  MAIN APP (CHAT INTERFACE)
# ══════════════════════════════════════════════════════════════════════════════
def show_main_app():
    username     = st.session_state.username
    info         = get_user_info(username)
    is_connected = check_ollama_connection()
    model_name   = get_model_name()

    # Grab any pending generation queued by sidebar buttons BEFORE rendering sidebar
    pending_prompt = st.session_state.pop("pending_prompt", None)
    pending_mode   = st.session_state.pop("pending_mode",   "generate")

    # ── Sidebar ───────────────────────────────────────────────────────────────
    with st.sidebar:
        try:
            with open("logo.png", "rb") as _f:
                _logo_b64 = base64.b64encode(_f.read()).decode()
            _logo_html = f"<img src='data:image/png;base64,{_logo_b64}' style='width:220px;margin-bottom:4px;'>"
        except FileNotFoundError:
            try:
                with open("logo.jpeg", "rb") as _f:
                    _logo_b64 = base64.b64encode(_f.read()).decode()
                _logo_html = f"<img src='data:image/jpeg;base64,{_logo_b64}' style='width:220px;margin-bottom:4px;'>"
            except FileNotFoundError:
                _logo_html = "<span style='font-size:1.8rem;'>📖</span>"
        st.markdown(
            f"<div style='text-align:center;padding:10px 0 4px;'>"
            f"{_logo_html}"
            f"</div>",
            unsafe_allow_html=True,
        )
        st.markdown(
            f"<div style='text-align:center;margin:6px 0;'>"
            f"<span class='user-badge'>👤 {username}</span></div>",
            unsafe_allow_html=True,
        )
        if is_connected:
            st.markdown(
                f"<div style='text-align:center;margin-bottom:4px;'>"
                f"<span class='status-online'>● ONLINE</span></div>",
                unsafe_allow_html=True,
            )
        else:
            st.markdown(
                "<div style='text-align:center;margin-bottom:4px;'>"
                "<span class='status-offline'>● OFFLINE — Start Ollama</span></div>",
                unsafe_allow_html=True,
            )

        st.markdown("---")
        st.markdown("<div class='sidebar-section-header'>Story Controls</div>", unsafe_allow_html=True)

        genre          = st.selectbox("Genre",  ["Romantic","Horror","Fantasy","Mystery","Humor","Suspense","Sci-Fi","Adventure"])
        writing_mode   = st.selectbox("Writing Mode", ["Beginning","Continue","Climax","Ending"])
        tone           = st.selectbox("Tone",   ["Dramatic","Lighthearted","Dark","Whimsical","Suspenseful","Melancholic","Epic"])
        creativity     = st.slider("Creativity", 0, 100, 50)
        response_length = st.selectbox("Response Length", ["Short","Medium","Long"], index=1)

        has_story = any(m["role"] == "assistant" for m in st.session_state.messages)

        if st.button("Continue Story", type="primary", use_container_width=True,
                     disabled=not is_connected or not has_story):
            st.session_state.pending_prompt = "Continue the story in an engaging and creative way."
            st.session_state.pending_mode   = "continue"
            st.rerun()

        st.markdown("---")
        st.markdown("<div class='sidebar-section-header'>Quick Story Prompts</div>", unsafe_allow_html=True)

        quick_prompts = [
            ("💕", "Start a romantic story"),
            ("👻", "Create a horror tale"),
            ("🗡️", "Write a fantasy adventure"),
            ("🔍", "Develop a mystery plot"),
            ("😄", "Craft a humorous scene"),
            ("⚡", "Build suspense"),
        ]
        for emoji, label in quick_prompts:
            if st.button(f"{emoji}  {label}", key=f"qp_{label}", use_container_width=True):
                st.session_state.pending_prompt = label
                st.session_state.pending_mode   = "generate"
                st.rerun()

        st.markdown("---")
        st.markdown("<div class='sidebar-section-header'>Story History</div>", unsafe_allow_html=True)

        recent = get_stories(username)[:5]
        if recent:
            for s in recent:
                _s_snippet = s['prompt'][:70] + "…" if len(s['prompt']) > 70 else s['prompt']
                st.markdown(
                    f"<div class='history-card'>"
                    f"<span class='msg-tag tag-genre'>{s['genre']}</span>"
                    f"<div class='history-prompt-label'>Prompt</div>"
                    f"<div class='history-title'>{_s_snippet}</div>"
                    f"<div class='history-date'>{format_date(s['created_at'])}</div>"
                    f"</div>",
                    unsafe_allow_html=True,
                )
            if st.button(f"View All {info['story_count']} Stories →", use_container_width=True, key="goto_history"):
                st.session_state.page = "history"
                st.rerun()
        else:
            st.markdown(
                "<div style='color:#555577;font-size:0.82rem;text-align:center;'>No stories yet!</div>",
                unsafe_allow_html=True,
            )

        st.markdown("---")
        c_new, c_out = st.columns(2)
        with c_new:
            if st.button("New Story", use_container_width=True):
                st.session_state.messages = []
                st.rerun()
        with c_out:
            if st.button("Logout", use_container_width=True):
                st.session_state.logged_in = False
                st.session_state.username  = ""
                st.session_state.messages  = []
                st.session_state.page      = "main"
                st.rerun()

    # Apply genre-specific animated background theme
    apply_genre_background(genre)

    # ── Main chat area ────────────────────────────────────────────────────────
    st.markdown("<div class='main-title'>STORY WAVER</div>", unsafe_allow_html=True)
    st.markdown("<div class='main-subtitle'>Begin Your Creative Story Journey ✨</div>",
                unsafe_allow_html=True)

    # ── Stats dashboard (always shown on welcome) ──
    if not st.session_state.messages and not pending_prompt:
        _total, _words, _fav = _get_user_stats(username)
        _words_str = f"{_words:,}" if _words >= 1000 else str(_words)
        st.markdown(
            f"<div class='stats-row'>"
            f"<div class='stat-card'><span class='stat-icon'>📚</span>"
            f"<span class='stat-value'>{_total}</span>"
            f"<span class='stat-label'>Stories Created</span></div>"
            f"<div class='stat-card'><span class='stat-icon'>✍️</span>"
            f"<span class='stat-value'>{_words_str}</span>"
            f"<span class='stat-label'>Words Written</span></div>"
            f"<div class='stat-card'><span class='stat-icon'>🎭</span>"
            f"<span class='stat-value' style='font-size:1.1rem;letter-spacing:-0.5px;'>{_fav}</span>"
            f"<span class='stat-label'>Favourite Genre</span></div>"
            f"</div>",
            unsafe_allow_html=True,
        )

    # ── Story phase progress bar (always visible) ──
    _phases      = ["Beginning", "Continue", "Climax", "Ending"]
    _phase_idx   = _phases.index(writing_mode) if writing_mode in _phases else 0
    _session_wc  = sum(len(m["content"].split()) for m in st.session_state.messages
                       if m["role"] == "assistant" and not m.get("blocked"))
    _nodes_html  = ""
    for idx, phase in enumerate(_phases):
        if idx < _phase_idx:
            dot_cls = "done"; lbl_cls = "done"
        elif idx == _phase_idx:
            dot_cls = "active"; lbl_cls = "active"
        else:
            dot_cls = ""; lbl_cls = ""
        _nodes_html += f"<div class='phase-node'><div class='phase-dot {dot_cls}'></div><div class='phase-label {lbl_cls}'>{phase}</div></div>"
        if idx < len(_phases) - 1:
            line_cls = "done" if idx < _phase_idx else ""
            _nodes_html += f"<div class='phase-line {line_cls}'></div>"
    _words_note = f"✍️ {_session_wc:,} words in current session" if _session_wc > 0 else "✍️ Start writing to track your word count"
    st.markdown(
        f"<div class='progress-wrap'>"
        f"<div class='progress-title'>Story Progress</div>"
        f"<div class='progress-phases'>{_nodes_html}</div>"
        f"<div class='progress-words'>{_words_note}</div>"
        f"</div>",
        unsafe_allow_html=True,
    )

    if not st.session_state.messages and not pending_prompt:
        st.markdown(
            "<div style='text-align:center;padding:40px 0 60px;'>"
            "<div style='font-size:0.85rem;color:#252535;letter-spacing:2px;"
            "text-transform:uppercase;font-weight:600;'>"
            "Type a prompt below to begin your story"
            "</div></div>",
            unsafe_allow_html=True,
        )

    # ── Render existing conversation ──
    for msg in st.session_state.messages:
        if msg["role"] == "user":
            with st.chat_message("user", avatar="👤"):
                st.markdown(msg["content"])
        else:
            with st.chat_message("assistant", avatar="📖"):
                if msg.get("blocked"):
                    st.markdown(msg["content"], unsafe_allow_html=True)
                else:
                    st.markdown(
                        f"<span class='msg-tag tag-genre'>{msg.get('genre','')}</span>"
                        f"<span class='msg-tag tag-mode'>{msg.get('writing_mode','')}</span>"
                        f"<span class='msg-tag tag-tone'>{msg.get('tone','')}</span>",
                        unsafe_allow_html=True,
                    )
                    st.markdown(msg["content"])

    # ── Generation logic ──────────────────────────────────────────────────────
    def run_generation(prompt_text: str, mode: str):
        # Resolve context (last assistant turn)
        context = ""
        if mode == "continue":
            for m in reversed(st.session_state.messages):
                if m["role"] == "assistant":
                    context = m["content"]
                    break

        effective_input = prompt_text.strip() or "Continue the story in an engaging way."
        effective_mode  = "Continue" if mode == "continue" else writing_mode

        # ── Guardrail check ──
        is_allowed, blocked_category = check_guardrails(effective_input)
        if not is_allowed:
            refusal_html = (
                "<div style='border-left:3px solid #cc4455;padding:14px 18px;"
                "background:#1a0a0f;border-radius:10px;color:#dd8899;'>"
                "<b style='font-size:1.05rem;'>⚠️ This is not related to our story.</b><br><br>"
                f"Your message contains <b>{blocked_category}</b>, which is offensive and "
                "falls outside the purpose of Story Waver. I am here exclusively to help you "
                "craft and explore creative stories — not to engage with harmful, offensive, "
                "or off-topic content.<br><br>"
                "<span style='color:#aa6677;font-size:0.92rem;'>Please keep our conversation "
                "respectful and story-focused. Try something like:<br>"
                "<i>\"Write a mystery story about a detective who discovers a secret...\"</i><br>"
                "or pick a Quick Story Prompt from the sidebar! ✍️</span>"
                "</div>"
            )
            st.session_state.messages.append({"role": "user", "content": effective_input})
            st.session_state.messages.append({"role": "assistant", "content": refusal_html, "blocked": True})
            st.rerun()
            return

        # Show user bubble
        with st.chat_message("user", avatar="👤"):
            st.markdown(effective_input)
        st.session_state.messages.append({"role": "user", "content": effective_input})

        # Stream assistant bubble
        with st.chat_message("assistant", avatar="📖"):
            st.markdown(
                f"<span class='msg-tag tag-genre'>{genre}</span>"
                f"<span class='msg-tag tag-mode'>{effective_mode}</span>"
                f"<span class='msg-tag tag-tone'>{tone}</span>",
                unsafe_allow_html=True,
            )
            full_text = st.write_stream(
                generate_story_stream(
                    user_input=effective_input,
                    genre=genre,
                    writing_mode=effective_mode,
                    tone=tone,
                    creativity=creativity,
                    response_length=response_length,
                    story_context=context,
                )
            )

        st.session_state.messages.append({
            "role": "assistant", "content": full_text,
            "genre": genre, "writing_mode": effective_mode, "tone": tone,
        })

        _clean_prompt = effective_input.strip().strip("?.!,;:-_/\\\"'")
        if full_text.strip() and len(_clean_prompt) >= 8:
            save_story(username=username, prompt=effective_input, story=full_text,
                       genre=genre, writing_mode=effective_mode, tone=tone)

        st.rerun()

    # Handle pending (sidebar buttons / quick prompts)
    if pending_prompt:
        run_generation(pending_prompt, pending_mode)

    # ── Chat input — auto-docked to bottom by Streamlit ──
    placeholder = "Describe your story idea…" if is_connected else "Ollama is offline — run: ollama serve"
    if prompt := st.chat_input(placeholder, disabled=not is_connected):
        run_generation(prompt, "generate")


# ══════════════════════════════════════════════════════════════════════════════
#  ROUTER
# ══════════════════════════════════════════════════════════════════════════════
if not st.session_state.logged_in:
    show_auth_page()
elif st.session_state.page == "history":
    show_history_page()
else:
    show_main_app()
